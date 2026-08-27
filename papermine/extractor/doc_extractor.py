"""本地文档抽取：纯文本、PDF、DOCX、PPTX 与中英文 OCR。"""
from __future__ import annotations

import hashlib
import io
import os
import threading
from dataclasses import dataclass, field
from typing import Any, Iterable, List, Optional, Set

from ..models import Asset

TEXT_EXTS = {".md", ".markdown", ".txt", ".rst", ".tex"}
DEFAULT_MAX_CHARS = 200_000
DEFAULT_MAX_PAGES = 100
DEFAULT_MAX_IMAGES = 60
MAX_IMAGE_PIXELS = 12_000_000
MIN_PDF_TEXT_CHARS = 40


@dataclass
class ExtractionResult:
    """一次文档抽取的文本及非致命诊断。"""

    text: str = ""
    warnings: List[str] = field(default_factory=list)
    truncated: bool = False


class _Collector:
    def __init__(self, max_chars: int) -> None:
        self.max_chars = max(0, max_chars)
        self.parts: List[str] = []
        self.length = 0
        self.truncated = False

    @property
    def full(self) -> bool:
        return self.length >= self.max_chars

    def add(self, value: Any) -> None:
        if self.full or value is None:
            if value:
                self.truncated = True
            return
        text = (
            str(value)
            .replace("\x00", "")
            .replace("\ufffe", "")
            .replace("\uffff", "")
            .replace("\u00ad", "")
            .strip()
        )
        if not text:
            return
        remaining = self.max_chars - self.length
        piece = text[:remaining]
        self.parts.append(piece)
        self.length += len(piece)
        if len(text) > remaining:
            self.truncated = True

    def render(self) -> str:
        return "\n".join(self.parts)


_OCR_ENGINE: Optional[Any] = None
_OCR_LOCK = threading.Lock()
_OCR_RUN_LOCK = threading.Lock()


def _get_ocr_engine() -> Any:
    global _OCR_ENGINE
    with _OCR_LOCK:
        if _OCR_ENGINE is None:
            from rapidocr_onnxruntime import RapidOCR

            _OCR_ENGINE = RapidOCR()
        return _OCR_ENGINE


def _ocr_image(image: Any) -> str:
    """对 PIL 图像执行本地中英文 OCR；结果按阅读顺序拼接。"""
    from PIL import Image
    import numpy as np

    if not isinstance(image, Image.Image):
        image = Image.open(io.BytesIO(image))
    image = image.convert("RGB")
    pixels = image.width * image.height
    if pixels > MAX_IMAGE_PIXELS:
        scale = (MAX_IMAGE_PIXELS / float(pixels)) ** 0.5
        image = image.resize(
            (max(1, int(image.width * scale)), max(1, int(image.height * scale)))
        )
    engine = _get_ocr_engine()
    with _OCR_RUN_LOCK:
        result, _ = engine(np.asarray(image))
    if not result:
        return ""
    return "\n".join(
        str(item[1]).strip()
        for item in result
        if isinstance(item, (list, tuple)) and len(item) >= 2 and str(item[1]).strip()
    )


def _ocr_blobs(
    blobs: Iterable[bytes], collector: _Collector, warnings: List[str], max_images: int
) -> None:
    seen: Set[bytes] = set()
    count = 0
    for blob in blobs:
        if collector.full:
            break
        signature = hashlib.sha256(blob).digest()
        if signature in seen:
            continue
        seen.add(signature)
        if count >= max_images:
            warnings.append("内嵌图片数量超过 OCR 预算，剩余图片未处理")
            collector.truncated = True
            break
        count += 1
        try:
            collector.add(_ocr_image(blob))
        except Exception as exc:
            warnings.append("第 {} 张内嵌图片 OCR 失败：{}".format(count, exc))


def _extract_plain(path: str, collector: _Collector) -> None:
    with open(path, "r", encoding="utf-8", errors="replace") as handle:
        collector.add(handle.read(collector.max_chars + 1))


def _extract_pdf(
    path: str, collector: _Collector, warnings: List[str], max_pages: int
) -> None:
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(path)
    try:
        page_count = len(pdf)
        if page_count > max_pages:
            warnings.append("PDF 共 {} 页，仅处理前 {} 页".format(page_count, max_pages))
            collector.truncated = True
        for index in range(min(page_count, max_pages)):
            if collector.full:
                break
            page = pdf[index]
            try:
                text_page = page.get_textpage()
                try:
                    text = text_page.get_text_range() or ""
                finally:
                    text_page.close()
                collector.add(text)
                if len("".join(text.split())) < MIN_PDF_TEXT_CHARS:
                    bitmap = page.render(scale=2.0)
                    try:
                        collector.add(_ocr_image(bitmap.to_pil()))
                    finally:
                        bitmap.close()
            except Exception as exc:
                warnings.append("PDF 第 {} 页解析失败：{}".format(index + 1, exc))
            finally:
                page.close()
    finally:
        pdf.close()


def _iter_docx_images(document: Any) -> Iterable[bytes]:
    for part in document.part.package.iter_parts():
        if str(getattr(part, "content_type", "")).startswith("image/"):
            yield part.blob


def _extract_docx(
    path: str, collector: _Collector, warnings: List[str], max_images: int
) -> None:
    from docx import Document

    document = Document(path)
    for paragraph in document.paragraphs:
        collector.add(paragraph.text)
    for table in document.tables:
        for row in table.rows:
            collector.add("\t".join(cell.text for cell in row.cells))
    for section in document.sections:
        for container in (section.header, section.footer):
            for paragraph in container.paragraphs:
                collector.add(paragraph.text)
            for table in container.tables:
                for row in table.rows:
                    collector.add("\t".join(cell.text for cell in row.cells))
    _ocr_blobs(_iter_docx_images(document), collector, warnings, max_images)


def _iter_shape_text(shape: Any) -> Iterable[str]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from _iter_shape_text(child)
        return
    if getattr(shape, "has_text_frame", False):
        yield shape.text
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            yield "\t".join(cell.text for cell in row.cells)


def _iter_pptx_images(presentation: Any) -> Iterable[bytes]:
    for part in presentation.part.package.iter_parts():
        if str(getattr(part, "content_type", "")).startswith("image/"):
            yield part.blob


def _extract_pptx(
    path: str,
    collector: _Collector,
    warnings: List[str],
    max_pages: int,
    max_images: int,
) -> None:
    from pptx import Presentation

    presentation = Presentation(path)
    slides = list(presentation.slides)
    if len(slides) > max_pages:
        warnings.append("PPTX 共 {} 页，仅处理前 {} 页".format(len(slides), max_pages))
        collector.truncated = True
    for index, slide in enumerate(slides[:max_pages]):
        collector.add("[幻灯片 {}]".format(index + 1))
        for shape in slide.shapes:
            for text in _iter_shape_text(shape):
                collector.add(text)
        try:
            if slide.has_notes_slide:
                collector.add("[演讲者备注]\n" + slide.notes_slide.notes_text_frame.text)
        except (AttributeError, ValueError):
            warnings.append("幻灯片 {} 的备注无法读取".format(index + 1))
    _ocr_blobs(_iter_pptx_images(presentation), collector, warnings, max_images)


def extract_asset_text(
    root: str,
    asset: Asset,
    max_chars: int = DEFAULT_MAX_CHARS,
    max_pages: int = DEFAULT_MAX_PAGES,
    max_images: int = DEFAULT_MAX_IMAGES,
) -> ExtractionResult:
    """读取一个文档资产；格式或依赖异常时返回诊断而非抛出。"""
    full = os.path.join(root, asset.path)
    ext = os.path.splitext(asset.path)[1].lower()
    collector = _Collector(max_chars)
    warnings: List[str] = []
    try:
        if ext in TEXT_EXTS or asset.kind == "readme":
            _extract_plain(full, collector)
        elif ext == ".pdf":
            _extract_pdf(full, collector, warnings, max_pages)
        elif ext == ".docx":
            _extract_docx(full, collector, warnings, max_images)
        elif ext == ".pptx":
            _extract_pptx(full, collector, warnings, max_pages, max_images)
    except ImportError as exc:
        warnings.append("缺少文档解析依赖：{}".format(exc))
    except Exception as exc:
        warnings.append("文档无法解析：{}".format(exc))
    return ExtractionResult(text=collector.render(), warnings=warnings, truncated=collector.truncated)


def read_asset_text(root: str, asset: Asset, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    """兼容既有接口，返回本地抽取出的纯文本。"""
    return extract_asset_text(root, asset, max_chars=max_chars).text
